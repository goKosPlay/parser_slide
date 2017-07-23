"""
Create By Kos
"""
from pyquery import PyQuery as pq
import urllib
import re
import reportlab.lib.utils
import reportlab.pdfgen.canvas
import sys
import argparse
import os
from pptx import Presentation
from pptx.util import Cm,Inches
import urllib.request
from PIL import Image


def speakerdeck_pptx(info, title):
    imageResult = re.findall(r'"original":"(.*?)"', str(info))
    prs = Presentation()

    for image in imageResult:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        imageFile = image.split('/')[-1]
        urllib.request.urlretrieve(image, imageFile)
        im = Image.open(imageFile)
        top = Inches(0)
        left = Inches(0)
        width, height = im.size
        width = Cm(width * 0.04 / 1.5)
        height = Cm(height * 0.04 / 1.5)
        pic = slide.shapes.add_picture(imageFile, left, top, width=width, height=height)
        im.close()
        os.remove(imageFile)
        print("--> %s" % image)
    prs.save(title.strip().replace(' ', '').replace("/", "_") + '.pptx')


def speakerdeck_pdf(info, title):
    imageResult = re.findall(r'"original":"(.*?)"', str(info))
    pdf = reportlab.pdfgen.canvas.Canvas(title.strip().replace(' ', '').replace("/", "_") + '.pdf')
    for image in imageResult:
        page = reportlab.lib.utils.ImageReader(image)
        page_width, page_height = page.getSize()
        pdf.setPageSize((page_width, page_height))
        pdf.drawImage(page, 0, 0, page_width, page_height)
        pdf.showPage()
        print("--> %s" % image)
    pdf.save()


def parse_speakerdeck(url, exportType='pdf'):
    try:
        htmlTag = pq(url)
        lastResult = htmlTag("div.speakerdeck-embed")
        for data in lastResult:
            info = pq("https://speakerdeck.com/player/%s" % data.attrib["data-id"])
            title = info('title').html()
            if exportType == 'pdf':
                speakerdeck_pdf(info, title)
            elif exportType == 'pptx':
                speakerdeck_pptx(info, title)
            else:
                speakerdeck_pdf(info, title)

    except urllib.error.HTTPError as err:
        if '404' in err:
            print("not found page.")
        else:
            print(err)


def parse_page(url, type):
    if 'speakerdeck.com' in url:
        parse_speakerdeck(url, type)


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Flip a switch by setting a flag")
    parser.add_argument('-type', help="匯出 pdf, 或者 pptx ,其中一個類型,預設沒有輸出為 pdf")
    parser.add_argument('-url', help="網址")
    args = parser.parse_args()
    if len(sys.argv) < 2:
        print('python parse_slide.py [My URL]')
    else:
        parse_page(args.url, args.type)
